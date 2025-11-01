#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
سكريبت لإنشاء عرض تقديمي PowerPoint جاهز
يتطلب تثبيت: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    # إنشاء عرض تقديمي جديد
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # الشريحة 1: العنوان الرئيسي
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # شريحة فارغة
    
    # خلفية زرقاء للشريحة الأولى
    background = slide1.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(41, 128, 185)
    
    # العنوان الرئيسي
    title_box = slide1.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "عرض تقديمي احترافي"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER
    
    # العنوان الفرعي
    subtitle_box = slide1.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(0.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "مُعد بواسطة Blackbox"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = RGBColor(255, 255, 255)
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # التاريخ
    date_box = slide1.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(0.5))
    date_frame = date_box.text_frame
    date_frame.text = "نوفمبر 2025"
    date_para = date_frame.paragraphs[0]
    date_para.font.size = Pt(18)
    date_para.font.color.rgb = RGBColor(236, 240, 241)
    date_para.alignment = PP_ALIGN.CENTER
    
    # الشريحة 2: جدول المحتويات
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # عنوان الشريحة
    title2_box = slide2.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(0.8))
    title2_frame = title2_box.text_frame
    title2_frame.text = "جدول المحتويات"
    title2_para = title2_frame.paragraphs[0]
    title2_para.font.size = Pt(40)
    title2_para.font.bold = True
    title2_para.font.color.rgb = RGBColor(41, 128, 185)
    title2_para.alignment = PP_ALIGN.RIGHT
    
    # المحتويات
    content_items = [
        "1. المقدمة",
        "2. الأهداف الرئيسية",
        "3. التحليل والبيانات",
        "4. الحلول المقترحة",
        "5. الخطة التنفيذية",
        "6. الخلاصة"
    ]
    
    content_box = slide2.shapes.add_textbox(Inches(2), Inches(2), Inches(6), Inches(4))
    content_frame = content_box.text_frame
    content_frame.text = content_items[0]
    
    for item in content_items[1:]:
        p = content_frame.add_paragraph()
        p.text = item
        p.space_before = Pt(15)
    
    for paragraph in content_frame.paragraphs:
        paragraph.font.size = Pt(28)
        paragraph.font.color.rgb = RGBColor(52, 73, 94)
        paragraph.alignment = PP_ALIGN.RIGHT
    
    # الشريحة 3: المقدمة
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # عنوان
    title3_box = slide3.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(0.8))
    title3_frame = title3_box.text_frame
    title3_frame.text = "المقدمة"
    title3_para = title3_frame.paragraphs[0]
    title3_para.font.size = Pt(40)
    title3_para.font.bold = True
    title3_para.font.color.rgb = RGBColor(41, 128, 185)
    title3_para.alignment = PP_ALIGN.RIGHT
    
    # المحتوى
    intro_text = """
    • نقدم لكم عرضاً شاملاً يغطي جميع الجوانب المهمة
    
    • تم إعداد هذا العرض بعناية فائقة لتلبية احتياجاتكم
    
    • يتضمن العرض تحليلاً دقيقاً وحلولاً عملية
    
    • نسعى لتحقيق أفضل النتائج من خلال هذه الخطة
    """
    
    intro_box = slide3.shapes.add_textbox(Inches(1.5), Inches(2), Inches(7), Inches(4))
    intro_frame = intro_box.text_frame
    intro_frame.text = intro_text.strip()
    intro_frame.word_wrap = True
    
    for paragraph in intro_frame.paragraphs:
        paragraph.font.size = Pt(22)
        paragraph.font.color.rgb = RGBColor(52, 73, 94)
        paragraph.alignment = PP_ALIGN.RIGHT
        paragraph.space_after = Pt(20)
    
    # الشريحة 4: الأهداف
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # عنوان
    title4_box = slide4.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(0.8))
    title4_frame = title4_box.text_frame
    title4_frame.text = "الأهداف الرئيسية"
    title4_para = title4_frame.paragraphs[0]
    title4_para.font.size = Pt(40)
    title4_para.font.bold = True
    title4_para.font.color.rgb = RGBColor(41, 128, 185)
    title4_para.alignment = PP_ALIGN.RIGHT
    
    # الأهداف في صناديق ملونة
    goals = [
        ("تحسين الأداء", RGBColor(46, 204, 113)),
        ("زيادة الكفاءة", RGBColor(52, 152, 219)),
        ("تطوير المهارات", RGBColor(155, 89, 182)),
        ("تحقيق النجاح", RGBColor(241, 196, 15))
    ]
    
    y_position = 2
    for goal, color in goals:
        box = slide4.shapes.add_textbox(Inches(2), Inches(y_position), Inches(6), Inches(0.8))
        frame = box.text_frame
        frame.text = goal
        para = frame.paragraphs[0]
        para.font.size = Pt(28)
        para.font.bold = True
        para.font.color.rgb = RGBColor(255, 255, 255)
        para.alignment = PP_ALIGN.CENTER
        
        # إضافة خلفية ملونة
        fill = box.fill
        fill.solid()
        fill.fore_color.rgb = color
        
        y_position += 1.1
    
    # الشريحة 5: البيانات والإحصائيات
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # عنوان
    title5_box = slide5.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(0.8))
    title5_frame = title5_box.text_frame
    title5_frame.text = "البيانات والإحصائيات"
    title5_para = title5_frame.paragraphs[0]
    title5_para.font.size = Pt(40)
    title5_para.font.bold = True
    title5_para.font.color.rgb = RGBColor(41, 128, 185)
    title5_para.alignment = PP_ALIGN.RIGHT
    
    # إحصائيات
    stats = [
        ("85%", "نسبة النجاح"),
        ("120+", "مشروع منجز"),
        ("50+", "عميل راضٍ"),
        ("24/7", "دعم متواصل")
    ]
    
    x_positions = [1, 3.5, 6, 8.5]
    for i, (stat, desc) in enumerate(stats):
        if i < len(x_positions):
            # الرقم
            stat_box = slide5.shapes.add_textbox(Inches(x_positions[i] - 0.5), Inches(2.5), Inches(1.5), Inches(0.8))
            stat_frame = stat_box.text_frame
            stat_frame.text = stat
            stat_para = stat_frame.paragraphs[0]
            stat_para.font.size = Pt(36)
            stat_para.font.bold = True
            stat_para.font.color.rgb = RGBColor(41, 128, 185)
            stat_para.alignment = PP_ALIGN.CENTER
            
            # الوصف
            desc_box = slide5.shapes.add_textbox(Inches(x_positions[i] - 0.5), Inches(3.5), Inches(1.5), Inches(0.5))
            desc_frame = desc_box.text_frame
            desc_frame.text = desc
            desc_para = desc_frame.paragraphs[0]
            desc_para.font.size = Pt(16)
            desc_para.font.color.rgb = RGBColor(52, 73, 94)
            desc_para.alignment = PP_ALIGN.CENTER
    
    # الشريحة 6: الخلاصة
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # خلفية خضراء
    background6 = slide6.background
    fill6 = background6.fill
    fill6.solid()
    fill6.fore_color.rgb = RGBColor(39, 174, 96)
    
    # العنوان
    conclusion_box = slide6.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    conclusion_frame = conclusion_box.text_frame
    conclusion_frame.text = "شكراً لحسن استماعكم"
    conclusion_para = conclusion_frame.paragraphs[0]
    conclusion_para.font.size = Pt(48)
    conclusion_para.font.bold = True
    conclusion_para.font.color.rgb = RGBColor(255, 255, 255)
    conclusion_para.alignment = PP_ALIGN.CENTER
    
    # معلومات التواصل
    contact_box = slide6.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(0.5))
    contact_frame = contact_box.text_frame
    contact_frame.text = "للتواصل: info@example.com"
    contact_para = contact_frame.paragraphs[0]
    contact_para.font.size = Pt(20)
    contact_para.font.color.rgb = RGBColor(255, 255, 255)
    contact_para.alignment = PP_ALIGN.CENTER
    
    # حفظ الملف
    prs.save('presentation.pptx')
    print("✅ تم إنشاء ملف PowerPoint بنجاح!")
    print("📁 اسم الملف: presentation.pptx")
    print("📊 عدد الشرائح: 6")

if __name__ == "__main__":
    try:
        create_presentation()
    except ImportError:
        print("❌ خطأ: يرجى تثبيت المكتبة المطلوبة")
        print("قم بتشغيل: pip install python-pptx")
    except Exception as e:
        print(f"❌ حدث خطأ: {e}")
